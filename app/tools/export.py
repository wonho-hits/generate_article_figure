"""Export pipeline: package session artifacts as downloadable files.

Three formats, each tied to its compatible session kind:
- SVG    → for SVG sessions (Path A): pass-through utf-8 encode.
- image  → for raster sessions (Path C): pass-through with detected MIME.
- PPTX   → for raster sessions: single-slide pptx with the image centered.
"""

from __future__ import annotations

import io
from dataclasses import dataclass

from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches

from app.tools.export_svg_pptx import build_svg_embedded_pptx
from app.tools.raster_illustration import detect_image_mime

PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


@dataclass(frozen=True)
class ExportResult:
    content: bytes
    media_type: str
    filename: str


def export_svg(svg_string: str, *, session_id: str) -> ExportResult:
    if not svg_string:
        raise ValueError("svg_string is empty")
    return ExportResult(
        content=svg_string.encode("utf-8"),
        media_type="image/svg+xml",
        filename=_filename(session_id, "svg"),
    )


def export_image(image_bytes: bytes, *, session_id: str) -> ExportResult:
    if not image_bytes:
        raise ValueError("image_bytes is empty")
    mime = detect_image_mime(image_bytes)
    if mime == "application/octet-stream":
        raise ValueError("session bytes are not a recognized image format")
    return ExportResult(
        content=image_bytes,
        media_type=mime,
        filename=_filename(session_id, _ext_for_mime(mime)),
    )


def export_pptx(image_bytes: bytes, *, session_id: str) -> ExportResult:
    if not image_bytes:
        raise ValueError("image_bytes is empty")
    if detect_image_mime(image_bytes) == "application/octet-stream":
        raise ValueError("session bytes are not a recognized image format")

    pptx_bytes = _build_single_image_pptx(image_bytes)
    return ExportResult(
        content=pptx_bytes,
        media_type=PPTX_MIME,
        filename=_filename(session_id, "pptx"),
    )


def export_pptx_from_svg(svg_string: str, *, session_id: str) -> ExportResult:
    """Build a PPTX with the SVG embedded as a vector image.

    Modern PowerPoint (2016+) renders the SVG natively. Users can
    right-click → "Convert to Shape" to break it into native editable
    shapes (text boxes, ellipses, paths).
    """
    if not svg_string or not svg_string.strip():
        raise ValueError("svg_string is empty")
    pptx_bytes = build_svg_embedded_pptx(svg_string)
    return ExportResult(
        content=pptx_bytes,
        media_type=PPTX_MIME,
        filename=_filename(session_id, "pptx"),
    )


def _build_single_image_pptx(image_bytes: bytes) -> bytes:
    """One blank slide with the image centered with a 0.5-inch margin.

    Aspect-fit: scaled down to fit the available area without distortion.
    """
    img_stream = io.BytesIO(image_bytes)
    pil = PILImage.open(img_stream)
    img_w, img_h = pil.size

    prs = Presentation()
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin = Inches(0.5)
    avail_w = slide_w - 2 * margin
    avail_h = slide_h - 2 * margin

    img_aspect = img_w / img_h
    box_aspect = avail_w / avail_h
    if img_aspect > box_aspect:
        final_w = avail_w
        final_h = int(avail_w / img_aspect)
    else:
        final_h = avail_h
        final_w = int(avail_h * img_aspect)

    left = (slide_w - final_w) // 2
    top = (slide_h - final_h) // 2

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    img_stream.seek(0)
    slide.shapes.add_picture(
        img_stream, left=left, top=top, width=final_w, height=final_h
    )

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _filename(session_id: str, ext: str) -> str:
    short = session_id[:8] if len(session_id) >= 8 else session_id
    # session_id is a UUID we generated — already safe for filenames.
    return f"figure_{short}.{ext}"


def _ext_for_mime(mime: str) -> str:
    return {
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/webp": "webp",
        "image/gif": "gif",
        "image/svg+xml": "svg",
    }.get(mime, "bin")
