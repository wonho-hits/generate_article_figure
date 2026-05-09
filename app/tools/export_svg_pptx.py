"""SVG-embedded PPTX export.

Produces a .pptx where the SVG is included as an OOXML SVG image. Modern
PowerPoint (2016+) renders it as a vector graphic and the user can
right-click → "Convert to Shape" to break it into native editable shapes.
Older PowerPoint sees only the placeholder PNG fallback.

This sidesteps python-pptx (which doesn't expose SVG embed) by post-processing
the resulting PPTX ZIP: add the SVG, its relationship, the content-type entry,
and the asvg:svgBlip extension on the existing picture's <a:blip>.
"""

from __future__ import annotations

import base64
import io
import zipfile

from lxml import etree
from pptx import Presentation
from pptx.util import Inches

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
RELS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
RT_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"

# 1×1 transparent PNG placeholder. PowerPoint needs *some* PNG fallback. Modern
# PowerPoint shows the SVG; this PNG is only seen by versions that don't speak
# SVG embed (PowerPoint 2013 and older).
_PLACEHOLDER_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAAC0lEQVR4nGNgYAAAAAMAAWgmWQ0AAAAASUVORK5CYII="
)


def build_svg_embedded_pptx(svg_string: str) -> bytes:
    """Build a one-slide PPTX with the SVG embedded as a vector image."""
    if not svg_string or not svg_string.strip():
        raise ValueError("svg_string is empty")

    width, height = _parse_svg_dimensions(svg_string)
    aspect = width / height if height > 0 else 4 / 3

    base_pptx = _build_baseline_pptx(_PLACEHOLDER_PNG, aspect)
    return _inject_svg(base_pptx, svg_string)


def _parse_svg_dimensions(svg_string: str) -> tuple[float, float]:
    """Pull (width, height) from viewBox or width/height attrs. Default 800×600."""
    try:
        root = etree.fromstring(svg_string.encode("utf-8"))
        vb = root.get("viewBox")
        if vb:
            parts = vb.split()
            if len(parts) == 4:
                return float(parts[2]), float(parts[3])
        w_attr = root.get("width")
        h_attr = root.get("height")
        if w_attr and h_attr:
            return _strip_unit(w_attr), _strip_unit(h_attr)
    except Exception:
        pass
    return 800.0, 600.0


def _strip_unit(s: str) -> float:
    digits = "".join(c for c in s if c.isdigit() or c == ".")
    return float(digits) if digits else 0.0


def _build_baseline_pptx(png_bytes: bytes, aspect: float) -> bytes:
    """One slide with a picture sized to the SVG aspect ratio, centered with margin."""
    prs = Presentation()
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin = Inches(0.5)
    avail_w = slide_w - 2 * margin
    avail_h = slide_h - 2 * margin

    box_aspect = avail_w / avail_h
    if aspect > box_aspect:
        final_w = avail_w
        final_h = int(avail_w / aspect)
    else:
        final_h = avail_h
        final_w = int(avail_h * aspect)

    left = (slide_w - final_w) // 2
    top = (slide_h - final_h) // 2

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(
        io.BytesIO(png_bytes), left=left, top=top, width=final_w, height=final_h
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _inject_svg(pptx_bytes: bytes, svg_string: str) -> bytes:
    """Modify a PPTX to add the SVG + extension references."""
    with zipfile.ZipFile(io.BytesIO(pptx_bytes), "r") as in_zip:
        members = {name: in_zip.read(name) for name in in_zip.namelist()}

    members["[Content_Types].xml"] = _add_svg_content_type(members["[Content_Types].xml"])

    rels_key = "ppt/slides/_rels/slide1.xml.rels"
    members[rels_key], svg_rid = _add_svg_relationship(members[rels_key])

    slide_key = "ppt/slides/slide1.xml"
    members[slide_key] = _inject_asvg_blip(members[slide_key], svg_rid)

    members["ppt/media/image1.svg"] = svg_string.encode("utf-8")

    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as out_zip:
        for name, data in members.items():
            out_zip.writestr(name, data)
    return out.getvalue()


def _add_svg_content_type(xml_bytes: bytes) -> bytes:
    root = etree.fromstring(xml_bytes)
    for default in root.findall(f"{{{CT_NS}}}Default"):
        if default.get("Extension") == "svg":
            return xml_bytes  # already present
    new_default = etree.SubElement(root, f"{{{CT_NS}}}Default")
    new_default.set("Extension", "svg")
    new_default.set("ContentType", "image/svg+xml")
    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def _add_svg_relationship(rels_xml: bytes) -> tuple[bytes, str]:
    root = etree.fromstring(rels_xml)
    ns = f"{{{RELS_NS}}}"
    existing_ids = []
    for rel in root.findall(f"{ns}Relationship"):
        rid = rel.get("Id", "")
        if rid.startswith("rId"):
            try:
                existing_ids.append(int(rid[3:]))
            except ValueError:
                pass
    next_id = max(existing_ids, default=0) + 1
    new_rid = f"rId{next_id}"

    new_rel = etree.SubElement(root, f"{ns}Relationship")
    new_rel.set("Id", new_rid)
    new_rel.set("Type", RT_IMAGE)
    new_rel.set("Target", "../media/image1.svg")
    return (
        etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True),
        new_rid,
    )


def _inject_asvg_blip(slide_xml: bytes, svg_rid: str) -> bytes:
    """Add asvg:svgBlip extension inside the picture's <a:blip>.

    PowerPoint requires the namespace prefix to be exactly 'asvg' — without
    it, the extension is silently ignored and the file falls back to the PNG.
    Using `nsmap={"asvg": ASVG_NS}` on the SubElement forces lxml to declare
    the namespace with that exact prefix at this element's scope.
    """
    root = etree.fromstring(slide_xml)
    blip = root.find(f".//{{{A_NS}}}blip")
    if blip is None:
        raise RuntimeError("no <a:blip> element in slide XML")

    extLst = blip.find(f"{{{A_NS}}}extLst")
    if extLst is None:
        extLst = etree.SubElement(blip, f"{{{A_NS}}}extLst")

    ext = etree.SubElement(extLst, f"{{{A_NS}}}ext")
    ext.set("uri", SVG_EXT_URI)
    svg_blip = etree.SubElement(
        ext,
        f"{{{ASVG_NS}}}svgBlip",
        nsmap={"asvg": ASVG_NS},
    )
    svg_blip.set(f"{{{R_NS}}}embed", svg_rid)

    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
